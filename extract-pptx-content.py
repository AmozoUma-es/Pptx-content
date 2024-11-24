import sys
from pptx import Presentation

def extract_text_from_ppt(ppt_file, output_file):
    """
    Extracts text from a PowerPoint file and writes it to a text file.
    Each slide's content is separated by '------------------'.
    """
    try:
        # Load the PowerPoint presentation
        presentation = Presentation(ppt_file)
        
        # Open the output file to write the content
        with open(output_file, 'w', encoding='utf-8') as file:
            for slide_index, slide in enumerate(presentation.slides):
                # Extract text from each text box in the slide
                slide_content = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            slide_content.append(paragraph.text)
                
                # Write the slide content to the file
                file.write(f"Slide {slide_index + 1}:\n")
                file.write("\n".join(slide_content))
                file.write("\n")
                file.write("-" * 18 + "\n")
        
        print(f"The content of the presentation has been saved to '{output_file}'.")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    # Check if the correct number of arguments is provided
    if len(sys.argv) != 3:
        print("Usage: python extract-pptx-content.py <pptx_file> <txt_file>")
        sys.exit(1)
    
    # Input and output file paths
    ppt_file = sys.argv[1]
    output_file = sys.argv[2]
    
    extract_text_from_ppt(ppt_file, output_file)
